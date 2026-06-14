import streamlit as st
import io
import os
import tempfile
import pandas as pd
from pypdf import PdfWriter, PdfReader
from pdf2docx import Converter
import fitz  # PyMuPDF
from PIL import Image
import chromadb
from chromadb.utils import embedding_functions
import requests
import hashlib

# ==========================================
# 1. UI THEME & PAGE CONFIGURATION
# ==========================================
st.set_page_config(page_title="Local PDF Studio", layout="wide", initial_sidebar_state="collapsed")

# Custom CSS for App Background, Tab Styling, and Card Layout
st.markdown("""
    <style>
    /* Change the main app background to a light grey */
    .stApp {
        background-color: #f5f5fa;
    }

    .main-header { font-size: 32px; font-weight: bold; text-align: center; color: #333333; margin-bottom: 5px; }
    .sub-header { text-align: center; color: #666666; margin-bottom: 30px; }

    /* Custom tab bar (st.radio acting as tabs) - CENTER ALIGNED */
    div[data-testid="stRadio"] {
        display: flex;
        justify-content: center;
        align-items: center;
        width: 100%;
    }
    div[data-testid="stRadio"] > div[role="radiogroup"] {
        gap: 8px;
        justify-content: center;
        align-items: center;
        flex-wrap: wrap;
        border-bottom: 1px solid #e6e6ec;
        padding-bottom: 0;
        margin-bottom: 20px;
        width: 100%;
    }
    div[data-testid="stRadio"] label {
        padding: 10px 20px;
        font-weight: 600;
        color: #555555;
        border-radius: 0;
        border-bottom: 2px solid transparent;
        margin: 0 !important;
        cursor: pointer;
        display: flex;
        justify-content: center;
        align-items: center;
        text-align: center;
    }
    div[data-testid="stRadio"] label p {
        text-align: center !important;
        margin: 0 !important;
        width: 100%;
    }
    div[data-testid="stRadio"] label:has(input:checked) {
        color: #E2574C;
        border-bottom: 2px solid #E2574C;
    }
    div[data-testid="stRadio"] label > div:first-child {
        display: none !important; /* hide the circular radio dot */
    }

    /* ===================== CARD GRID ===================== */
    /* The bordered container Streamlit renders for st.container(border=True) */
    div[data-testid="stVerticalBlockBorderWrapper"]:has(div.tool-card-marker) {
        background-color: #ffffff;
        border: 1px solid #e6e6ec !important;
        border-radius: 10px !important;
        box-shadow: 0 2px 6px rgba(0,0,0,0.04);
        transition: all 0.2s ease-in-out;
        padding: 4px;
        margin-bottom: 4px;
    }
    div[data-testid="stVerticalBlockBorderWrapper"]:has(div.tool-card-marker):hover {
        border-color: #E2574C !important;
        box-shadow: 0 6px 14px rgba(226, 87, 76, 0.12);
        transform: translateY(-2px);
    }

    /* Icon badge */
    .tool-icon {
        width: 42px;
        height: 42px;
        border-radius: 9px;
        background-color: #E2574C;
        color: #ffffff;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 22px;
        margin-bottom: 10px;
    }

    /* Tool title button - looks like a heading, not a button */
    div.tool-card-marker + div div[data-testid="stButton"] button {
        background-color: transparent !important;
        border: none !important;
        box-shadow: none !important;
        padding: 0 !important;
        height: auto !important;
        width: auto !important;
        justify-content: flex-start !important;
        text-align: left !important;
        font-size: 17px !important;
        font-weight: 600 !important;
        color: #2b2b2b !important;
    }
    div.tool-card-marker + div div[data-testid="stButton"] button p {
        font-size: 17px !important;
        font-weight: 600 !important;
        text-align: left !important;
        margin: 0 !important;
    }
    div.tool-card-marker + div div[data-testid="stButton"] button:hover p {
        color: #E2574C !important;
    }
    div.tool-card-marker + div div[data-testid="stButton"] button:hover {
        background-color: transparent !important;
    }

    /* Description text under the title */
    .tool-desc {
        color: #777777;
        font-size: 13.5px;
        line-height: 1.4;
        margin-top: 6px;
    }
    /* ===================== END CARD GRID ===================== */

    /* Action Button Styling (The red execute buttons inside the tools) */
    .action-btn div[data-testid="stButton"] button {
        height: auto;
        width: auto;
        background-color: #E2574C !important;
        color: white !important;
        text-align: center;
        font-weight: bold !important;
        border-radius: 8px;
        padding: 10px 24px !important;
        border: none !important;
        box-shadow: none !important;
        justify-content: center !important;
    }
    .action-btn div[data-testid="stButton"] button p {
        font-weight: bold !important;
        color: white !important;
        font-size: 15px !important;
        text-align: center !important;
    }
    .action-btn div[data-testid="stButton"] button:hover { background-color: #c9443a !important; border-color: #c9443a !important; }

    /* Back button styling */
    div[data-testid="stButton"]:has(button:contains("Back")) button { width: auto !important; }
    </style>
""", unsafe_allow_html=True)
# ==========================================
# RAG CONFIG & HELPER FUNCTIONS (AI Summarizer)
# ==========================================
LM_STUDIO_URL = "http://localhost:1234/v1/chat/completions"
LM_STUDIO_MODEL = "local-model"
COLLECTION_NAME = "pdf_chunks"

@st.cache_resource
def get_chroma_collection():
    client = chromadb.PersistentClient(path="./chroma_db")
    embed_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name="all-MiniLM-L6-v2"
    )
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME, embedding_function=embed_fn
    )
    return collection

def extract_pages_from_pdf_rag(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        pages.append((i + 1, page.get_text()))
    return pages

def split_text(text, chunk_size=800, overlap=100):
    text = text.strip()
    if len(text) <= chunk_size:
        return [text] if text else []
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + chunk_size
        if end >= text_len:
            chunks.append(text[start:].strip())
            break
        boundary = -1
        for sep in [". ", "! ", "? ", "\n\n", "\n"]:
            idx = text.rfind(sep, start, end)
            if idx > boundary:
                boundary = idx + len(sep)
        if boundary <= start:
            idx = text.rfind(" ", start, end)
            boundary = idx if idx > start else end
        chunk = text[start:boundary].strip()
        if chunk:
            chunks.append(chunk)
        start = max(boundary - overlap, start + 1)
    return chunks

def chunk_pages(pages, chunk_size=800, overlap=100):
    chunks = []
    for page_num, text in pages:
        if not text.strip():
            continue
        for piece in split_text(text, chunk_size, overlap):
            if piece:
                chunks.append((piece, page_num))
    return chunks

def add_chunks_to_db(chunks, source_name):
    collection = get_chroma_collection()
    ids, documents, metadatas = [], [], []
    for i, (text, page_num) in enumerate(chunks):
        uid = hashlib.md5(f"{source_name}_p{page_num}_{i}".encode()).hexdigest()
        ids.append(uid)
        documents.append(text)
        metadatas.append({"source": source_name, "page": page_num})
    collection.add(documents=documents, ids=ids, metadatas=metadatas)

def retrieve_context(query, k=4):
    collection = get_chroma_collection()
    results = collection.query(query_texts=[query], n_results=k)
    docs = results["documents"][0]
    metas = results["metadatas"][0]
    return list(zip(docs, metas))

def is_summary_request(query):
    keywords = ["summarize", "summary", "overview", "main points", "key points", "tl;dr", "gist"]
    return any(k in query.lower() for k in keywords)

def get_all_chunks_for_source(source_name):
    collection = get_chroma_collection()
    results = collection.get(where={"source": source_name})
    docs = results["documents"]
    metas = results["metadatas"]
    if not docs:
        return []
    paired = sorted(zip(docs, metas), key=lambda x: x[1]["page"])
    return [d for d, m in paired]

def call_llm(prompt):
    payload = {
        "model": LM_STUDIO_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
    }
    resp = requests.post(LM_STUDIO_URL, json=payload, timeout=300)
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]

def summarize_pdf_rag(source_name):
    chunks = get_all_chunks_for_source(source_name)
    if not chunks:
        return "No document found. Please upload and process a PDF first."
    mini_summaries = []
    progress = st.progress(0, text="Summarizing sections...")
    for idx, chunk in enumerate(chunks):
        prompt = f"Summarize the following text in 2-3 sentences:\n\n{chunk}"
        mini_summaries.append(call_llm(prompt))
        progress.progress((idx + 1) / len(chunks), text=f"Summarized {idx+1}/{len(chunks)} sections")
    progress.empty()
    combined = "\n".join(mini_summaries)
    final_prompt = f"Combine these section summaries into one coherent overall summary:\n\n{combined}"
    return call_llm(final_prompt)

def ask_llm(question, retrieved):
    context_blocks = []
    for text, meta in retrieved:
        context_blocks.append(f"[Source: {meta['source']} | Page: {meta['page']}]\n{text}")
    context = "\n\n".join(context_blocks)
    prompt = f"""Answer the question using only the context below. Each context block is labeled with its source PDF and page number.
After your answer, cite the sources you used in the format (filename, page X).
If the answer isn't in the context, say you don't know.

Context:
{context}

Question: {question}
Answer:"""
    return call_llm(prompt)

# ==========================================
# 2. STATE MANAGEMENT & TOOL DATA
# ==========================================
if "active_tool" not in st.session_state:
    st.session_state.active_tool = None

TAB_NAMES = ["All", "Organize PDF", "Optimize PDF", "Convert PDF", "Edit PDF", "PDF Security", "PDF Intelligence"]

if "active_tab" not in st.session_state:
    st.session_state.active_tab = TAB_NAMES[0]
if "rag_messages" not in st.session_state:
    st.session_state.rag_messages = []
if "rag_last_pdf" not in st.session_state:
    st.session_state.rag_last_pdf = None
def set_tool(tool_name, tab_name):
    st.session_state.active_tool = tool_name
    st.session_state.active_tab = tab_name

def go_back():
    st.session_state.active_tool = None

TOOL_CATEGORIES = {
    "Organize PDF": [
        {"name": "Merge PDF", "icon": "🔗", "desc": "Combine PDFs in the order you want with the easiest PDF merger available."},
        {"name": "Split PDF", "icon": "✂️", "desc": "Separate one page or a whole set for easy conversion into independent PDF files."}
    ],
    "Optimize PDF": [
        {"name": "Compress PDF", "icon": "📦", "desc": "Reduce file size while optimizing for maximal PDF quality.                       	"},
        {"name": "Repair PDF", "icon": "🛠️", "desc": "Repair a corrupted PDF and recover data from damaged files.            "},
        {"name": "OCR PDF", "icon": "🔎", "desc": "Easily convert scanned PDFs into searchable, selectable documents.        "}
    ],
    "Convert PDF": [
        {"name": "PDF to Word", "icon": "📄", "desc": "Convert your PDFs into easily editable Word documents (.docx).        "},
        {"name": "PDF to PowerPoint", "icon": "📊", "desc": "Turn your PDF files into easy to edit PowerPoint slideshows.    "},
        {"name": "PDF to Excel", "icon": "📈", "desc": "Pull data straight from PDFs into Excel spreadsheets in seconds.     "},
        {"name": "Word to PDF", "icon": "📝", "desc": "Make Word documents easy to read by converting them to PDF.           "},
        {"name": "PowerPoint to PDF", "icon": "🖥️", "desc": "Make PowerPoint presentations easy to view by converting to PDF."},
        {"name": "Excel to PDF", "icon": "🧮", "desc": "Make Excel spreadsheets easy to read by converting them to PDF.      "},
        {"name": "PDF to JPG", "icon": "🖼️", "desc": "Convert each PDF page into a JPG or extract all images contained in a PDF."},
        {"name": "JPG to PDF", "icon": "🌄", "desc": "Convert JPG images to PDF in seconds, adjusting orientation and margins."},
        {"name": "HTML to PDF", "icon": "🌐", "desc": "Convert webpages in HTML to PDF, copy and paste a URL or upload a file."},
        {"name": "PDF to PDF/A", "icon": "🗄️", "desc": "Transform your PDF to PDF/A, the ISO-standardized version for long-term archiving."}
    ],
    "Edit PDF": [
        {"name": "Edit PDF", "icon": "✏️", "desc": "Add text, images, shapes or annotations to a PDF document.                 "},
        {"name": "Watermark", "icon": "💧", "desc": "Stamp an image or text over your PDF in seconds.                          "},
        {"name": "Rotate PDF", "icon": "🔄", "desc": "Rotate multiple PDFs at once, choosing the orientation you need.         "},
        {"name": "Page numbers", "icon": "🔢", "desc": "Add page numbers into PDFs with ease, choosing position and style.     "},
        {"name": "Crop PDF", "icon": "🔳", "desc": "Crop margins of PDF documents or select specific areas to remove.          "},
        {"name": "PDF Forms", "icon": "🧾", "desc": "Add fillable fields like text boxes, checkboxes, and signatures.          "}
    ],
    "PDF Security": [
        {"name": "Sign PDF", "icon": "✍️", "desc": "Sign yourself or request electronic signatures from others.                "},
        {"name": "Unlock PDF", "icon": "🔓", "desc": "Remove PDF password security, giving you freedom to use your PDFs.       "},
        {"name": "Protect PDF", "icon": "🔒", "desc": "Protect PDF files with a password to prevent unauthorized access.       "},
        {"name": "Compare PDF", "icon": "🆚", "desc": "Compare two PDF files and easily spot the differences between them.     "},
        {"name": "Redact PDF", "icon": "⬛", "desc": "Permanently remove sensitive information from a PDF document.            "}
    ],
    "PDF Intelligence": [
        {"name": "AI Summarizer", "icon": "🤖", "desc": "Summarize long PDF documents instantly using a local or API-based AI model."}
    ]
}

def render_tool_grid(tool_list, tab_prefix, tab_name):
    cols = st.columns(4, gap="medium")
    for index, tool in enumerate(tool_list):
        with cols[index % 4]:
            with st.container(border=True):
                st.markdown('<div class="tool-card-marker"></div>', unsafe_allow_html=True)
                st.markdown(f'<div class="tool-icon">{tool.get("icon", "📄")}</div>', unsafe_allow_html=True)
                st.button(tool['name'], key=f"btn_{tab_prefix}_{tool['name']}", on_click=set_tool, args=(tool['name'], tab_name))
                st.markdown(f'<div class="tool-desc">{tool.get("desc", "")}</div>', unsafe_allow_html=True)

# ==========================================
# 3. MAIN APP ROUTING (UI VS LOGIC)
# ==========================================
st.markdown('<div class="main-header">Every tool you need to work with PDFs locally</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Merge, split, compress, convert, rotate, unlock and watermark PDFs securely.</div>', unsafe_allow_html=True)

if st.session_state.active_tool:
    current_tool = st.session_state.active_tool
    st.button(f"⬅️ Back to {st.session_state.active_tab}", on_click=go_back, key="btn_back")
    st.markdown("---")
    st.markdown(f"### {current_tool}")

    st.markdown('<div class="action-btn">', unsafe_allow_html=True)

    # ---------------- ORGANIZE PDF ----------------
    if current_tool == "Merge PDF":
        files = st.file_uploader("Select PDF files", accept_multiple_files=True, type=['pdf'])
        if st.button("Merge PDFs") and files:
            merger = PdfWriter()
            for f in files: merger.append(f)
            out = io.BytesIO()
            merger.write(out)
            st.success("PDFs merged successfully!")
            st.download_button("Download Merged PDF", out.getvalue(), "merged_document.pdf")

    elif current_tool == "Split PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        page = st.number_input("Extract Page Number", min_value=1, step=1)
        if st.button("Split PDF") and f:
            reader = PdfReader(f)
            if page <= len(reader.pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[page-1])
                out = io.BytesIO()
                writer.write(out)
                st.success(f"Page {page} extracted!")
                st.download_button("Download Extracted Page", out.getvalue(), f"page_{page}.pdf")
            else:
                st.error("Page number exceeds document length.")

    # ---------------- OPTIMIZE PDF ----------------
    elif current_tool == "Compress PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Compress PDF") and f:
            reader = PdfReader(f)
            writer = PdfWriter()
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            out = io.BytesIO()
            writer.write(out)
            st.success("Compression applied!")
            st.download_button("Download Compressed PDF", out.getvalue(), "compressed_document.pdf")

    elif current_tool == "Repair PDF":
        f = st.file_uploader("Select corrupted PDF", type=['pdf'])
        if st.button("Repair PDF") and f:
            try:
                doc = fitz.open(stream=f.read(), filetype="pdf")
                out_bytes = doc.tobytes(garbage=4, deflate=True)
                st.success("PDF Repaired!")
                st.download_button("Download Repaired PDF", out_bytes, "repaired.pdf")
            except Exception as e:
                st.error(f"Failed to repair: {e}")

    elif current_tool == "OCR PDF":
        st.warning("Requires Tesseract-OCR installed on your system.")
        f = st.file_uploader("Select Scanned PDF", type=['pdf'])
        if st.button("Run OCR") and f:
            with st.spinner("Processing OCR..."):
                try:
                    from pdf2image import convert_from_bytes
                    import pytesseract
                    from fpdf import FPDF
                    pages = convert_from_bytes(f.read(), dpi=300)
                    pdf = FPDF()
                    for page in pages:
                        text = pytesseract.image_to_string(page)
                        pdf.add_page()
                        pdf.set_font("Arial", size=12)
                        pdf.multi_cell(0, 10, text.encode('latin-1', 'ignore').decode('latin-1'))
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                        pdf.output(tmp.name)
                        with open(tmp.name, "rb") as out_f:
                            st.download_button("Download Searchable PDF", out_f.read(), "ocr_result.pdf")
                except Exception as e:
                    st.error(f"OCR Error: {e}")

    # ---------------- CONVERT PDF ----------------
    elif current_tool == "PDF to Word":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Convert to Word") and f:
            with open("temp.pdf", "wb") as tmp: tmp.write(f.getvalue())
            cv = Converter("temp.pdf")
            cv.convert("temp.docx")
            cv.close()
            with open("temp.docx", "rb") as word_file:
                st.success("Converted successfully!")
                st.download_button("Download Word Document", word_file.read(), "converted.docx")

    elif current_tool == "PDF to PowerPoint":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Convert to PPT") and f:
            with st.spinner("Converting to PPTX..."):
                try:
                    from pptx import Presentation
                    doc = fitz.open(stream=f.read(), filetype="pdf")
                    prs = Presentation()
                    blank_slide = prs.slide_layouts[6]
                    for i in range(len(doc)):
                        page = doc.load_page(i)
                        pix = page.get_pixmap(dpi=150)
                        img_path = f"temp_{i}.png"
                        pix.save(img_path)
                        slide = prs.slides.add_slide(blank_slide)
                        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                        os.remove(img_path)
                    out = io.BytesIO()
                    prs.save(out)
                    st.success("Converted to PowerPoint!")
                    st.download_button("Download PPT", out.getvalue(), "converted.pptx")
                except Exception as e:
                    st.error(f"Conversion failed: {e}")

    elif current_tool == "PDF to Excel":
        st.warning("Requires Ghostscript installed on your system.")
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Extract Tables to Excel") and f:
            with st.spinner("Extracting tables..."):
                try:
                    import camelot
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                        tmp.write(f.read())
                        tmp_path = tmp.name
                    tables = camelot.read_pdf(tmp_path, pages='all', flavor='stream')
                    if tables.n > 0:
                        out_path = tmp_path.replace(".pdf", ".xlsx")
                        tables.export(out_path, f='excel')
                        with open(out_path, "rb") as out_f:
                            st.success(f"Exported {tables.n} tables!")
                            st.download_button("Download Excel", out_f.read(), "extracted_data.xlsx")
                    else:
                        st.warning("No tables found.")
                except Exception as e:
                    st.error(f"Excel extraction failed: {e}")

    elif current_tool == "Word to PDF":
        st.warning("Requires Microsoft Word installed on Windows/Mac.")
        f = st.file_uploader("Select DOCX file", type=['docx'])
        if st.button("Convert to PDF") and f:
            try:
                from docx2pdf import convert as docx_convert
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp_in:
                    tmp_in.write(f.read())
                    in_path = tmp_in.name
                out_path = in_path.replace(".docx", ".pdf")
                docx_convert(in_path, out_path)
                with open(out_path, "rb") as out_f:
                    st.success("Converted Word to PDF!")
                    st.download_button("Download PDF", out_f.read(), "converted.pdf")
            except Exception as e:
                st.error(f"Conversion failed: {e}")

    elif current_tool in ["PowerPoint to PDF", "Excel to PDF"]:
        st.warning("Requires Windows OS and MS Office Installed.")
        f = st.file_uploader(f"Select {'PPTX' if 'PowerPoint' in current_tool else 'XLSX'} file")
        if st.button("Convert to PDF") and f:
            try:
                import comtypes.client
                ext = ".pptx" if "PowerPoint" in current_tool else ".xlsx"
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_in:
                    tmp_in.write(f.read())
                    in_path = os.path.abspath(tmp_in.name)
                out_path = os.path.abspath(tmp_in.name.replace(ext, ".pdf"))
                
                if "PowerPoint" in current_tool:
                    app = comtypes.client.CreateObject("Powerpoint.Application")
                    deck = app.Presentations.Open(in_path)
                    deck.SaveAs(out_path, 32)
                    deck.Close()
                else:
                    app = comtypes.client.CreateObject("Excel.Application")
                    wb = app.Workbooks.Open(in_path)
                    wb.ExportAsFixedFormat(0, out_path)
                    wb.Close()
                app.Quit()
                with open(out_path, "rb") as out_f:
                    st.success("Conversion complete!")
                    st.download_button("Download PDF", out_f.read(), "converted.pdf")
            except Exception as e:
                st.error(f"Conversion failed: {e}")

    elif current_tool == "PDF to JPG":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Convert to JPG") and f:
            doc = fitz.open(stream=f.read(), filetype="pdf")
            for i in range(min(len(doc), 3)):
                page = doc.load_page(i)
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                st.image(img_data, caption=f"Page {i+1}", width=300)
                st.download_button(f"Download Page {i+1}", img_data, f"page_{i+1}.png", key=f"dl_{i}")

    elif current_tool == "JPG to PDF":
        files = st.file_uploader("Select Images", accept_multiple_files=True, type=['jpg', 'jpeg', 'png'])
        if st.button("Convert to PDF") and files:
            images = [Image.open(io.BytesIO(img.read())).convert("RGB") for img in files]
            if images:
                out = io.BytesIO()
                images[0].save(out, format="PDF", save_all=True, append_images=images[1:])
                st.success("Images combined into PDF!")
                st.download_button("Download PDF", out.getvalue(), "images.pdf")

    elif current_tool == "HTML to PDF":
        url = st.text_input("Enter Website URL (e.g., https://example.com)")
        if st.button("Convert URL to PDF") and url:
            try:
                import pdfkit
                out_bytes = pdfkit.from_url(url, False)
                st.success("Website converted!")
                st.download_button("Download Website PDF", out_bytes, "website.pdf")
            except Exception as e:
                st.error(f"Conversion failed (Requires wkhtmltopdf installed): {e}")

    elif current_tool == "PDF to PDF/A":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Convert to PDF/A") and f:
            try:
                import ghostscript
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_in:
                    tmp_in.write(f.read())
                    in_path = tmp_in.name
                out_path = in_path.replace(".pdf", "_A.pdf")
                args = ["pdf2pdfa", "-dPDFA=1", "-dBATCH", "-dNOPAUSE", "-sProcessColorModel=DeviceCMYK", "-sDEVICE=pdfwrite", f"-sOutputFile={out_path}", in_path]
                ghostscript.Ghostscript(*args)
                with open(out_path, "rb") as out_f:
                    st.success("Converted to Archival Standard!")
                    st.download_button("Download PDF/A", out_f.read(), "archive.pdf")
            except Exception as e:
                st.error(f"Conversion failed: {e}")

    # ---------------- EDIT PDF ----------------
    elif current_tool == "Edit PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        text = st.text_input("Text to add")
        col1, col2, col3 = st.columns(3)
        page_num = col1.number_input("Page Num", min_value=1, step=1)
        x_cord = col2.number_input("X Coordinate", value=100)
        y_cord = col3.number_input("Y Coordinate", value=100)
        if st.button("Add Text") and f and text:
            doc = fitz.open(stream=f.read(), filetype="pdf")
            if page_num <= len(doc):
                page = doc[page_num - 1]
                page.insert_text(fitz.Point(x_cord, y_cord), text, fontsize=12, color=(1, 0, 0))
                st.success("Text Added!")
                st.download_button("Download Edited PDF", doc.tobytes(), "edited.pdf")

    elif current_tool == "Watermark":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        wm_text = st.text_input("Watermark Text", value="CONFIDENTIAL")
        if st.button("Apply Watermark") and f and wm_text:
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.setFont("Helvetica", 40)
                can.setFillGray(0.5, 0.5)
                can.saveState()
                can.translate(300, 400)
                can.rotate(45)
                can.drawCentredString(0, 0, wm_text)
                can.restoreState()
                can.save()
                packet.seek(0)
                watermark_pdf = PdfReader(packet)
                
                reader = PdfReader(f)
                writer = PdfWriter()
                for page in reader.pages:
                    page.merge_page(watermark_pdf.pages[0])
                    writer.add_page(page)
                out = io.BytesIO()
                writer.write(out)
                st.success("Watermark Applied!")
                st.download_button("Download Watermarked PDF", out.getvalue(), "watermarked.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

    elif current_tool == "Rotate PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        angle = st.selectbox("Rotation Angle", [90, 180, 270])
        if st.button("Rotate PDF") and f:
            reader = PdfReader(f)
            writer = PdfWriter()
            for page in reader.pages:
                page.rotate(angle)
                writer.add_page(page)
            out = io.BytesIO()
            writer.write(out)
            st.success("PDF Rotated!")
            st.download_button("Download Rotated PDF", out.getvalue(), "rotated.pdf")

    elif current_tool == "Page numbers":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        if st.button("Add Page Numbers") and f:
            try:
                from reportlab.pdfgen import canvas
                reader = PdfReader(f)
                writer = PdfWriter()
                for i, page in enumerate(reader.pages):
                    packet = io.BytesIO()
                    can = canvas.Canvas(packet)
                    can.drawString(500, 30, f"Page {i+1} of {len(reader.pages)}")
                    can.save()
                    packet.seek(0)
                    number_pdf = PdfReader(packet)
                    page.merge_page(number_pdf.pages[0])
                    writer.add_page(page)
                out = io.BytesIO()
                writer.write(out)
                st.success("Page Numbers Added!")
                st.download_button("Download Numbered PDF", out.getvalue(), "numbered.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

    elif current_tool == "Crop PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        st.write("Enter crop coordinates (in points):")
        c1, c2, c3, c4 = st.columns(4)
        left = c1.number_input("Left", value=50)
        bottom = c2.number_input("Bottom", value=50)
        right = c3.number_input("Right", value=550)
        top = c4.number_input("Top", value=750)
        if st.button("Crop PDF") and f:
            reader = PdfReader(f)
            writer = PdfWriter()
            for page in reader.pages:
                page.cropbox.lower_left = (left, bottom)
                page.cropbox.upper_right = (right, top)
                writer.add_page(page)
            out = io.BytesIO()
            writer.write(out)
            st.success("PDF Cropped!")
            st.download_button("Download Cropped PDF", out.getvalue(), "cropped.pdf")

    elif current_tool == "PDF Forms":
        st.info("Form Filling Demo")
        f = st.file_uploader("Select Fillable PDF", type=['pdf'])
        k = st.text_input("Field Name (e.g. 'Name')")
        v = st.text_input("Field Value (e.g. 'Om Bombate')")
        if st.button("Fill Form") and f and k and v:
            reader = PdfReader(f)
            writer = PdfWriter()
            writer.append_pages_from_reader(reader)
            writer.update_page_form_field_values(writer.pages[0], {k: v})
            out = io.BytesIO()
            writer.write(out)
            st.success("Form Filled!")
            st.download_button("Download Filled Form", out.getvalue(), "filled.pdf")

    # ---------------- PDF SECURITY ----------------
    elif current_tool == "Sign PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        cert = st.file_uploader("Select .pfx Certificate", type=['pfx'])
        pwd = st.text_input("Certificate Password", type="password")
        if st.button("Digitally Sign PDF") and f and cert and pwd:
            try:
                from pyhanko.sign import signers
                from pyhanko.pdf_utils.writer import copy_into_new_writer
                from pyhanko.pdf_utils.reader import PdfFileReader
                signer = signers.SimpleSigner.load_pkcs12(cert.read(), pwd.encode('utf-8'))
                pdf_reader = PdfFileReader(io.BytesIO(f.read()))
                pdf_writer = copy_into_new_writer(pdf_reader)
                out = io.BytesIO()
                signers.sign_pdf(pdf_writer, signers.PdfSignatureMetadata(field_name='Signature1'), signer=signer, in_place=False, out=out)
                st.success("PDF Digitally Signed!")
                st.download_button("Download Signed PDF", out.getvalue(), "signed.pdf")
            except Exception as e:
                st.error(f"Signing failed: {e}")

    elif current_tool == "Unlock PDF":
        f = st.file_uploader("Select Locked PDF file", type=['pdf'])
        pwd = st.text_input("Original Password", type="password")
        if st.button("Unlock PDF") and f and pwd:
            reader = PdfReader(f)
            if reader.is_encrypted:
                if reader.decrypt(pwd):
                    writer = PdfWriter()
                    for page in reader.pages: writer.add_page(page)
                    out = io.BytesIO()
                    writer.write(out)
                    st.success("Password Removed!")
                    st.download_button("Download Unlocked PDF", out.getvalue(), "unlocked.pdf")
                else:
                    st.error("Incorrect Password.")
            else:
                st.info("This PDF is not encrypted.")

    elif current_tool == "Protect PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        pwd = st.text_input("Type a secure password", type="password")
        if st.button("Protect PDF") and f and pwd:
            reader = PdfReader(f)
            writer = PdfWriter()
            for page in reader.pages: writer.add_page(page)
            writer.encrypt(pwd)
            out = io.BytesIO()
            writer.write(out)
            st.success("Document Encrypted securely!")
            st.download_button("Download Protected PDF", out.getvalue(), "protected.pdf")

    elif current_tool == "Compare PDF":
        f1 = st.file_uploader("Select Original PDF", type=['pdf'], key='f1')
        f2 = st.file_uploader("Select Modified PDF", type=['pdf'], key='f2')
        if st.button("Compare") and f1 and f2:
            try:
                from PIL import ImageChops
                doc1 = fitz.open(stream=f1.read(), filetype="pdf")
                doc2 = fitz.open(stream=f2.read(), filetype="pdf")
                pix1 = doc1[0].get_pixmap()
                pix2 = doc2[0].get_pixmap()
                img1 = Image.frombytes("RGB", [pix1.width, pix1.height], pix1.samples)
                img2 = Image.frombytes("RGB", [pix2.width, pix2.height], pix2.samples)
                diff = ImageChops.difference(img1, img2)
                if diff.getbbox():
                    st.warning("Differences found on Page 1!")
                    out = io.BytesIO()
                    diff.save(out, format="PNG")
                    st.image(out.getvalue(), caption="Visual Difference Map")
                else:
                    st.success("Page 1 is visually identical.")
            except Exception as e:
                st.error(f"Compare failed: {e}")

    elif current_tool == "Redact PDF":
        f = st.file_uploader("Select PDF file", type=['pdf'])
        word = st.text_input("Sensitive Word to Redact")
        if st.button("Redact Text") and f and word:
            doc = fitz.open(stream=f.read(), filetype="pdf")
            count = 0
            for page in doc:
                instances = page.search_for(word)
                for inst in instances:
                    page.add_redact_annot(inst, fill=(0, 0, 0))
                    count += 1
                page.apply_redactions()
            if count > 0:
                st.success(f"Redacted {count} instances of '{word}'.")
                st.download_button("Download Redacted PDF", doc.tobytes(garbage=4, deflate=True), "redacted.pdf")
            else:
                st.info("Word not found.")

    elif current_tool == "AI Summarizer":
        st.markdown("### Ask AI about your PDF")
        st.info("⚡ Powered by your local LM Studio (100% Offline & Secure)")

        f = st.file_uploader("Upload PDF Document", type=['pdf'])

        if f is not None:
            if st.button("Process PDF"):
                with st.spinner("Extracting pages and chunking..."):
                    file_bytes = f.read()
                    pages = extract_pages_from_pdf_rag(file_bytes)
                    chunks = chunk_pages(pages)
                    add_chunks_to_db(chunks, f.name)
                    st.session_state.rag_last_pdf = f.name
                st.success(f"Stored {len(chunks)} chunks from '{f.name}' ({len(pages)} pages) in vector DB.")

        st.divider()

        for msg in st.session_state.rag_messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if query := st.chat_input("Ask a question about your PDF...", key="rag_chat_input"):
            st.session_state.rag_messages.append({"role": "user", "content": query})
            with st.chat_message("user"):
                st.markdown(query)

            with st.spinner("Thinking..."):
                try:
                    if is_summary_request(query):
                        answer = summarize_pdf_rag(st.session_state.rag_last_pdf)
                        retrieved = []
                    else:
                        retrieved = retrieve_context(query)
                        if not retrieved:
                            answer = "No documents in the database yet. Please upload and process a PDF first."
                        else:
                            answer = ask_llm(query, retrieved)
                except requests.exceptions.ConnectionError:
                    answer = "⚠️ Could not connect to LM Studio. Make sure the local server is running on http://localhost:1234"
                    retrieved = []
                except Exception as e:
                    answer = f"⚠️ Error: {e}"
                    retrieved = []

            st.session_state.rag_messages.append({"role": "assistant", "content": answer})
            with st.chat_message("assistant"):
                st.markdown(answer)
                if retrieved:
                    st.markdown("**🔍 Chunks sent to LLM as context:**")
                    for i, (text, meta) in enumerate(retrieved):
                        label = f"Chunk {i+1} — {meta['source']} (Page {meta['page']})"
                        st.text_area(label, text, height=120, key=f"rag_chunk_{len(st.session_state.rag_messages)}_{i}")

    else:
        st.info("🚧 This specific tool interface is under construction. Check back later!")

    st.markdown('</div>', unsafe_allow_html=True)

# SCENARIO B: NO TOOL SELECTED, SHOW THE MENU GRIDS
else:
    # Custom tab bar (st.radio styled as tabs) so the selected tab can be
    # restored programmatically when returning from a tool.
    selected_tab = st.radio(
        "",
        TAB_NAMES,
        key="active_tab",
        horizontal=True,
        label_visibility="collapsed",
    )

    if selected_tab == "All":
        all_tools = []
        for category_tools in TOOL_CATEGORIES.values():
            all_tools.extend(category_tools)
        render_tool_grid(all_tools, "All", "All")
    elif selected_tab == "Organize PDF":
        render_tool_grid(TOOL_CATEGORIES["Organize PDF"], "Org", "Organize PDF")
    elif selected_tab == "Optimize PDF":
        render_tool_grid(TOOL_CATEGORIES["Optimize PDF"], "Opt", "Optimize PDF")
    elif selected_tab == "Convert PDF":
        render_tool_grid(TOOL_CATEGORIES["Convert PDF"], "Conv", "Convert PDF")
    elif selected_tab == "Edit PDF":
        render_tool_grid(TOOL_CATEGORIES["Edit PDF"], "Edit", "Edit PDF")
    elif selected_tab == "PDF Security":
        render_tool_grid(TOOL_CATEGORIES["PDF Security"], "Sec", "PDF Security")
    elif selected_tab == "PDF Intelligence":
        render_tool_grid(TOOL_CATEGORIES["PDF Intelligence"], "Intel", "PDF Intelligence")
